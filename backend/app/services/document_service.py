import os
import re
import json
import logging
import requests
from typing import List, Dict, Any, Optional
from openai import OpenAI

logger = logging.getLogger("jam_analyzer")

# Conditional imports for document parsers
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

from app.config import settings

class DocumentService:
    def __init__(self):
        self.api_key = settings.OPENAI_KEY or settings.OPENAI_API_KEY or os.getenv("OPENAI_KEY") or os.getenv("OPENAI_API_KEY")
        if self.api_key:
            self.client = OpenAI(api_key=self.api_key)
        else:
            self.client = None
            logger.warning("OPENAI_KEY/OPENAI_API_KEY environment variable is not configured for DocumentService.")

    def extract_text(self, file_path: str, filename: str) -> str:
        """
        Extracts raw text content from the uploaded file based on its extension.
        Logs detailed extraction status.
        """
        ext = os.path.splitext(filename.lower())[1]
        text_content = ""

        # Step 1: Diagnose input file details
        try:
            file_size_bytes = os.path.getsize(file_path)
            file_size_kb = round(file_size_bytes / 1024, 2)
        except Exception:
            file_size_kb = 0.0

        logger.info("--- Document Extraction Diagnostics ---")
        logger.info(f"File Name: {filename}")
        logger.info(f"File Type: {ext}")
        logger.info(f"File Size: {file_size_kb} KB")
        logger.info(f"Storage Path: {file_path}")

        try:
            if ext == ".pdf":
                if pypdf:
                    reader = pypdf.PdfReader(file_path)
                    text_content = "\n".join([page.extract_text() or "" for page in reader.pages])
                else:
                    logger.warning("pypdf is not installed. Using basic byte decoding fallback.")
                    text_content = self._fallback_text_extraction(file_path)
            
            elif ext in [".docx", ".doc"]:
                if docx and ext == ".docx":
                    doc = docx.Document(file_path)
                    text_content = "\n".join([p.text for p in doc.paragraphs])
                else:
                    logger.warning("docx parser fallback triggered.")
                    text_content = self._fallback_text_extraction(file_path)
            
            elif ext in [".pptx", ".ppt"]:
                if pptx and ext == ".pptx":
                    prs = pptx.Presentation(file_path)
                    slides_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slides_text.append(shape.text)
                    text_content = "\n".join(slides_text)
                else:
                    logger.warning("pptx parser fallback triggered.")
                    text_content = self._fallback_text_extraction(file_path)
            
            elif ext in [".txt", ".md"]:
                # Try standard UTF-8 first
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        text_content = f.read()
                except UnicodeDecodeError:
                    # Try Latin-1 fallback
                    logger.warning(f"UTF-8 decoding failed for {filename}. Trying Latin-1 fallback.")
                    try:
                        with open(file_path, "r", encoding="latin-1") as f:
                            text_content = f.read()
                    except Exception as e:
                        raise ValueError(f"Unsupported Encoding: Failed to read text file. Details: {e}")
            
            else:
                # Catch-all basic reading
                text_content = self._fallback_text_extraction(file_path)

        except Exception as e:
            logger.error(f"Error during file text extraction: {e}")
            raise ValueError(f"Text Extraction Failed: {str(e)}")

        # Cleanup text formatting slightly
        text_content = re.sub(r'\s+', ' ', text_content).strip()
        
        # Step 2: Validate extracted character count
        text_len = len(text_content)
        if text_len == 0:
            raise ValueError(f"No Content Extracted: The file '{filename}' was resolved as empty.")

        preview = text_content[:200] + "..." if text_len > 200 else text_content
        logger.info(f"Extracted Characters: {text_len}")
        logger.info(f"Extracted Text Preview: \"{preview}\"")
        logger.info("---------------------------------------")
            
        return text_content

    def _fallback_text_extraction(self, file_path: str) -> str:
        """
        Reads binary files extracting printable ASCII strings as fallback.
        """
        try:
            with open(file_path, "rb") as f:
                data = f.read(500000) # Limit to first 500KB for safety
                ascii_strings = re.findall(b"[a-zA-Z0-9\s\.,;:!@#\$%\^&\*\(\)\-_\+=\[\]\{\}<>\?\/\\|~`']{4,}", data)
                decoded = []
                for s in ascii_strings:
                    try:
                        decoded.append(s.decode("utf-8"))
                    except:
                        pass
                return "\n".join(decoded)
        except Exception as e:
            logger.error(f"Fallback text extraction failed: {e}")
            return ""

    def _call_llm_json(self, prompt: str, system_prompt: str = "You are an elite coach. Output clean JSON only.") -> Dict[str, Any]:
        """
        Unified LLM orchestrator using key fallback: OpenAI -> Gemini -> Groq.
        Logs details for diagnostics.
        """
        logger.info("--- AI LLM Request Diagnostics ---")
        logger.info(f"Document/Prompt Length: {len(prompt)} characters")
        logger.info(f"Prompt Sent: {prompt}")

        # Attempt 1: OpenAI
        if self.client:
            try:
                logger.info("Attempting OpenAI (model: gpt-4o) call...")
                res = self.client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": prompt}
                    ],
                    response_format={"type": "json_object"},
                    temperature=0.2
                )
                raw_response = res.choices[0].message.content.strip()
                logger.info(f"OpenAI Response Received: {raw_response}")
                logger.info("----------------------------------")
                return json.loads(raw_response)
            except Exception as e:
                logger.warning(f"OpenAI call failed: {e}")

        # Attempt 2: Gemini
        if settings.GEMINI_API_KEY:
            try:
                logger.info("Attempting Gemini (model: gemini-2.5-flash) call...")
                from google import genai
                from google.genai import types
                
                client = genai.Client(api_key=settings.GEMINI_API_KEY)
                response = client.models.generate_content(
                    model='gemini-2.5-flash',
                    contents=prompt,
                    config=types.GenerateContentConfig(
                        response_mime_type="application/json",
                        temperature=0.2,
                        system_instruction=system_prompt
                    )
                )
                if response.text:
                    raw_response = response.text.strip()
                    logger.info(f"Gemini Response Received: {raw_response}")
                    logger.info("----------------------------------")
                    return json.loads(raw_response)
            except Exception as e:
                logger.warning(f"Gemini call failed: {e}")

        # Attempt 3: Groq
        if settings.GROQ_API_KEY:
            try:
                logger.info("Attempting Groq (model: llama-3.3-70b-versatile) call...")
                # Note: Groq json mode requires that the prompt has the word "json"
                groq_prompt = prompt
                if "json" not in groq_prompt.lower():
                    groq_prompt += "\nReturn your response in clean JSON format."
                
                headers = {
                    "Authorization": f"Bearer {settings.GROQ_API_KEY}",
                    "Content-Type": "application/json"
                }
                payload = {
                    "model": "llama-3.3-70b-versatile",
                    "messages": [
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": groq_prompt}
                    ],
                    "response_format": {"type": "json_object"},
                    "temperature": 0.2
                }
                res = requests.post(
                    "https://api.groq.com/openai/v1/chat/completions",
                    headers=headers,
                    json=payload,
                    timeout=25
                )
                if res.status_code == 200:
                    raw_response = res.json()["choices"][0]["message"]["content"].strip()
                    logger.info(f"Groq Response Received: {raw_response}")
                    logger.info("----------------------------------")
                    return json.loads(raw_response)
                else:
                    logger.warning(f"Groq call failed with status code {res.status_code}: {res.text}")
            except Exception as e:
                logger.warning(f"Groq call failed: {e}")

        raise ValueError("AI Request Failed: No AI services (OpenAI, Gemini, Groq) were available or succeeded to evaluate the request.")

    def analyze_document_content(self, text: str) -> Dict[str, Any]:
        """
        AI reads the extracted text and outputs title, summary, key_concepts, keywords, and classification.
        No hardcoded fallbacks.
        """
        if not text.strip():
            raise ValueError("Empty Document: Content to analyze is empty.")

        truncated_text = text[:15000] # Limit tokens sent for summary
        
        prompt = f"""
        Analyze the following extracted document text and return a JSON object containing:
        1. "title": A descriptive title of the document.
        2. "summary": A concise 2-3 sentence summary of the document.
        3. "key_concepts": A list of the 3-5 most important conceptual themes/models in the text.
        4. "keywords": A list of 4-6 keyword tags.
        5. "learning_objectives": A list of 3-5 key learning objectives or takeaways from this document.
        6. "is_project_report": Boolean, set to true if the document presents a software development, hardware system, architecture, or academic project report.
        7. "is_resume": Boolean, set to true if the document is a CV, resume, or professional profile.

        Extracted Document Text:
        \"\"\"{truncated_text}\"\"\"

        Return ONLY the raw JSON object.
        """

        try:
            return self._call_llm_json(prompt, "You are a professional document parser and analyst. Output clean JSON only.")
        except Exception as e:
            logger.error(f"Failed document analysis: {e}")
            raise ValueError(f"AI Request Failed: Failed to analyze document content. Details: {str(e)}")

    def generate_topics(self, doc_analysis: Dict[str, Any], text: str) -> List[Dict[str, Any]]:
        """
        Generates 5-10 unique speaking topics from the document based on its category and contents.
        No hardcoded fallbacks.
        """
        title = doc_analysis.get("title", "Document")
        summary = doc_analysis.get("summary", "")
        concepts = ", ".join(doc_analysis.get("key_concepts", []))
        keywords = ", ".join(doc_analysis.get("keywords", []))

        # Determine type
        doc_type = "standard"
        if doc_analysis.get("is_project_report"):
            doc_type = "project_report"
        elif doc_analysis.get("is_resume"):
            doc_type = "resume"

        prompt = f"""
        You are a topic generation agent for Document-Based Presentation training.
        Based on this document:
        Title: {title}
        Summary: {summary}
        Key Concepts: {concepts}
        Keywords: {keywords}
        Document Classification: {doc_type.upper()}
        Extracted content preview:
        \"\"\"{text[:4000]}\"\"\"

        Generate a list of 5 to 10 unique, relevant speaking presentation topics strictly based on the extracted content.
        Ensure you match the Document Classification requirements:
        - If STANDARD: Generate conceptual topics mapping the document's themes. (Beginner, Intermediate, Advanced).
        - If PROJECT_REPORT: Automatically generate specific topics:
          1. Project Presentation Topic
          2. Architecture Design and Choices
          3. Technical Implementation Details
          4. Challenges & Future Improvements
        - If RESUME: Automatically generate professional topics:
          1. Professional Experience & Role Highlights
          2. Key Projects & Technical Achievements
          3. Career Goals & Fit
          4. Behavioral strengths based on resume

        Each topic in the list MUST be a JSON object with:
        - "topic": The speaking topic title.
        - "category": The category name (e.g. Technology, Engineering, Career).
        - "difficulty": "Beginner", "Intermediate", or "Advanced".
        - "talking_points": A list of 3 bullet talking points.
        - "estimated_speaking_time": Time in seconds (usually 60 or 90).
        - "keywords": 3 keywords associated with the topic.

        Return a JSON object containing a "topics" array of these topic objects.
        """

        try:
            res_data = self._call_llm_json(prompt, "You are a communication training coordinator. Output clean JSON only.")
            topics = res_data.get("topics", [])
            if not topics:
                raise ValueError("AI returned empty topics array.")
            return topics
        except Exception as e:
            logger.error(f"Failed to generate topics: {e}")
            raise ValueError(f"AI Request Failed: Failed to generate topics. Details: {str(e)}")

    def evaluate_speech_against_document(
        self, document_text: str, topic_title: str, transcript: str, communication_scores: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Performs comparative AI analysis between document content and user speaking transcript.
        Evaluates Content Accuracy, Concept Coverage, Topic Relevance, Understanding, Explanation Quality, and Tech Correctness.
        """
        truncated_doc = document_text[:12000] # Limit token inputs
        
        prompt = f"""
        You are an elite academic presentation evaluator and document presentation coach.
        Compare the user's spoken transcript against the source document to evaluate their content understanding.

        Topic User Chose to Present: "{topic_title}"
        User spoken transcript:
        \"\"\"{transcript}\"\"\"

        Source Document text:
        \"\"\"{truncated_doc}\"\"\"

        Evaluate and return a JSON object with:
        1. "accuracy_score": Integer (0-100), how factual and correct the user's statements are compared to the document.
        2. "coverage_score": Integer (0-100), what percentage of key concepts relevant to the topic they managed to mention.
        3. "understanding_score": Integer (0-100), overall metric of subject comprehension.
        4. "explanation_quality": Integer (0-100), evaluates structural clarity, flow, and educational clarity.
        5. "technical_correctness": Integer (0-100), score based on accuracy of technical vocabulary, algorithms, or facts from the text.
        6. "relevance_score": Integer (0-100), evaluation of whether the user stayed on-topic relative to the chosen topic title.
        7. "knowledge_gaps": List of objects representing concepts present in the document relevant to this topic that the user completely missed. Each object should have:
           - "concept": Name of the missing concept/topic.
           - "description": 1 sentence describing what was missed and where to find it.
        8. "suggested_improvements": List of 2-3 specific suggestions to improve their conceptual explanations.
        9. "coach_recommendations": List of 2-3 study prompts or coaching advice.
        10. "follow_up_questions": List of 3 dynamic follow-up questions testing content depth based on the document.

        Return ONLY the raw JSON object.
        """

        try:
            return self._call_llm_json(prompt, "You are a master presentation examiner. Output clean JSON only.")
        except Exception as e:
            logger.error(f"Failed comparative speech-to-document evaluation: {e}")
            raise ValueError(f"AI Request Failed: Failed to compare speech against document. Details: {str(e)}")

    def generate_viva_questions(self, document_text: str, mode: str) -> List[str]:
        """
        Generates 5 oral examination / interview questions based on document text and chosen mode.
        """
        truncated_doc = document_text[:12000]
        
        prompt = f"""
        You are an expert interviewer and academic coordinator.
        Based on the following document content:
        \"\"\"{truncated_doc}\"\"\"

        Generate exactly 5 challenging oral questions matching the mode: {mode.upper()}
        - VIVA: Focus on theory, validity, methodology, academic defense.
        - PROJECT: Focus on system design, database schemas, code flow, scaling, testing.
        - RESUME: Focus on experience highlights, candidate skills verification, technical questions about projects, HR behavioral fit.

        Return a JSON object containing a "questions" list of strings.
        """

        try:
            res_data = self._call_llm_json(prompt, "You are a professional panel examiner. Output clean JSON only.")
            questions = res_data.get("questions", [])
            if not questions or len(questions) < 5:
                raise ValueError(f"AI returned invalid questions list: {questions}")
            return questions
        except Exception as e:
            logger.error(f"Failed to generate viva questions: {e}")
            raise ValueError(f"AI Request Failed: Failed to generate viva questions. Details: {str(e)}")

    def evaluate_viva_response(self, document_text: str, question: str, answer_transcript: str) -> Dict[str, Any]:
        """
        Grades a single viva oral answer out of 100 and provides specific qualitative feedback.
        """
        truncated_doc = document_text[:12000]
        
        prompt = f"""
        Evaluate this candidate's oral response to a Viva question based on the source document:
        
        Question: "{question}"
        Candidate Spoken Response:
        \"\"\"{answer_transcript}\"\"\"

        Source Document:
        \"\"\"{truncated_doc}\"\"\"

        Generate a JSON object containing:
        1. "score": Integer (0-100), grading the correctness and depth of the answer relative to the document.
        2. "feedback": A 2-sentence summary feedback.
        3. "correct_elements": A list of points the candidate explained correctly.
        4. "incorrect_elements": A list of points the candidate got wrong or described weakly.
        5. "ideal_response": A sample expert-level response to the question.

        Return ONLY the raw JSON object.
        """

        try:
            return self._call_llm_json(prompt, "You are an oral board assessor. Output clean JSON only.")
        except Exception as e:
            logger.error(f"Failed to evaluate viva response: {e}")
            raise ValueError(f"AI Request Failed: Failed to grade viva answer. Details: {str(e)}")
